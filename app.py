from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from collections import defaultdict
from string import punctuation
from heapq import nlargest
from newspaper import Article
import nltk
import random
from textblob import TextBlob
from rake_nltk import Rake
from googleapiclient.discovery import build
import pprint
from nltk.corpus import movie_reviews
from pptx import Presentation
from pptx.util import Inches
from flask import json
from flask import request
from flask import Flask, Response
from flask_cors import CORS

class FrequencySummarizer:
    '''
    FrequencySummarizer Class used to find the maximum repetition of words in a particular context
    and reduce the repetition for better summary

    
    '''
    def __init__(self, min_freq=0.1, max_freq=0.9)-> None:
        ''' Initialize class Frequency Summarizer. Words below and above the max_freq and min_freq
        threshold will not be included by default '''
        self._min_freq = min_freq
        self._max_freq = max_freq
        self.end_word = set(stopwords.words('english') + list(punctuation))
    
    def quantify_freq(self, word_l)-> defaultdict:
        ''' Quantify frequency of all words in the text. Sentences are already been tokenized using
        nltk.tokenzie
        
        Parameters : word_l
        word_l -> an array of all sentences already been tokenized
        '''
        list_freq = defaultdict(int)
        for sentence in word_l:
            for word in sentence:
                if word not in self.end_word:
                    list_freq[word] += 1
        max_word_freq = float(max(list_freq.values()))
        for word in list_freq.keys():
            list_freq[word] = list_freq[word]/max_word_freq
            if (list_freq[word] >= self._max_freq) or (list_freq[word] <= self._min_freq):
                list_freq[word] = -1
        new_list_freq = defaultdict(int)
        for word in list_freq.keys():
            if list_freq[word] != -1:
                new_list_freq[word] = list_freq[word]
        return new_list_freq
        
    def summarize_text(self, t, n):
        '''Return a list of n sentences that represents the slide and the image.'''
        list_sentences = sent_tokenize(t)
        assert n <= len(list_sentences)
        word_sentences = [word_tokenize(s.lower()) for s in list_sentences]
        self._freq = self.quantify_freq(word_sentences)
        ranking = defaultdict(int)
        for i, sentence in enumerate(word_sentences):
            for w in sentence:
                ranking[i] += self._freq[w]
        sents_idx = self.ranks(ranking, n)
        return [list_sentences[j] for j in sents_idx]
    
    def ranks(self, ranked, n):
        ''' Returns the first n sentences with highest rank'''
        return nlargest(n, ranked, key=ranked.get)
   
# Google API : AIzaSyBTfn4y7zrQx1vUM88tjs1dC7wetpbmDJs
class Analyze:
    '''class analyze would analyze the text and prepare the right type of slides'''
    def __init__(self, text=""):
        self._text = text
        self._sentiment_analysis = 0.0
        self._sentiment_analysis_subjectivity = 0.0
        self._ranked_words = []
        self._place_interest = []
        self._selected_images = []
        self._selected_images_url = []
    
    
    def sentiment_analysis(self)->int:
        new_analysis = TextBlob(self._text)
        '''Subjective is personal feelings vs objective is pure facts'''
        self._sentiment_analysis = new_analysis.sentiment.polarity
        self._sentiment_analysis_subjectivity = new_analysis.subjectivity
        return self._sentiment_analysis
    
    def place_interest(self) -> list:
        '''Add dictionary element to this'''
        new_text = self._text
        new_list = dict()
        for sent in nltk.sent_tokenize(new_text):
            for chunk in nltk.ne_chunk(nltk.pos_tag(nltk.word_tokenize(sent))):
                if hasattr(chunk, 'label'):
                    if (chunk.label() not in new_list):
                        new_list[chunk.label()] = ' '.join(c[0] for c in chunk).lower()
                    else:
                        temporary = [new_list[chunk.label()]]
                        temporary.append((' '.join(c[0] for c in chunk)).lower())
                        new_list[chunk.label()] = temporary
        self._place_interest = new_list
        print(self._place_interest)
    
    def keyword_extraction(self):
        new_keyword_find = Rake()
        new_keyword_find.extract_keywords_from_text(self._text)
        ranked_words = new_keyword_find.get_ranked_phrases_with_scores()
        self._ranked_words = ranked_words        
    
    def google_search(self, **kwargs):
        my_api_key = "AIzaSyBTfn4y7zrQx1vUM88tjs1dC7wetpbmDJs"
        my_cse_id = "013938504458615515541:cmbncicbmdc"
        service = build('customsearch', "v1", developerKey=my_api_key)
        for val in self._selected_images:
            res = service.cse().list(q=val, cx=my_cse_id, searchType='image', num = 1, fileType='png', safe='off').execute()
            if not 'items' in res:
                print ('No result!!\nres is: {}'.format(res))
            else:
                for item in res['items']:
                    print('{}:\n\t{}'.format(item['title'],item['link']))
                    self._selected_images_url.append(item['link'])
                    
    
    def image_list(self):
        temporary = []
        temporary_ranked = []
        for val in self._place_interest.values():
            if type(val) == list:
                for i in val:
                    temporary.append(i)
            else:
                temporary.append(val)
        
        for val in self._ranked_words:
            temporary_ranked.append(val[1])
        self._selected_images.append(temporary_ranked[0])
        
        for val_temp in temporary:
            for val_temp_ranked in temporary_ranked:
                try:
                    if val_temp in val_temp_ranked:
                        self._selected_images.append(val_temp_ranked)
                except TypeError:
                    print("ok")

class Generate_Presentation:
    
    def __init__(self, images = [], summary = ""):
        self.images = images
        self.summary = summary
        
    def make_slide(self, title = "Title Not Specified"):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left,top,width,height)
        tf = txBox.text_frame
        tf.text = self.summary
        left = Inches(1)
        top = Inches(5)
        print(self.images)
        pic = slide.shapes.add_picture("test.png", left, top)
        prs.save("text.pptx")
        

app = Flask(__name__)
CORS(app)

@app.route('/', methods = ['GET', 'POST'])
def run_app():
    data = request.form.to_dict()
    text = data['value']
    new_analyze = Analyze(text)
    #new_analyze._sentiment_analysis()
    new_analyze.place_interest()
    new_analyze.keyword_extraction()
    new_analyze.image_list()
    new_analyze.google_search()
    new_summarize = Article("http://www.bbc.com/news/world-asia-43933332")
    new_summarize.download()
    new_summarize.parse()
    new_summarize.text = text
    new_summarize.nlp()
    summarize = new_summarize.summary
    
    return([summarize, new_analyze._selected_images_url])

if __name__ == "__main__":
    text = "Hurricane Carmen was the most intense tropical cyclone of the 1974 Atlantic hurricane season. A destructive and widespread storm, Carmen originated as a tropical disturbance that traveled westward from Africa, spawning a tropical depression east of the Lesser Antilles on August 29. Moving through the Caribbean Sea, it quickly strengthened to a Category 4 hurricane on the Saffir–Simpson Hurricane Scale, and made landfall on the Yucatán Peninsula. It turned north into the Gulf of Mexico, re-intensified, and made a second landfall in the marshland of southern Louisiana, dissipating over eastern Texas on September 10. Tropical cyclone watches and warnings had been issued for the storm, and around 100,000 residents left their homes and sought shelter. Damage was lighter than first feared, but the sugar industry suffered substantial losses. The hurricane killed 8 people and caused damage valued at $162 million. The name Carmen was retired from the list of Atlantic tropical cyclone names in 1975."
    new_analyze = Analyze(text)
    #new_analyze._sentiment_analysis()
    new_analyze.place_interest()
    new_analyze.keyword_extraction()
    new_analyze.image_list()
    print(new_analyze._ranked_words)
    print(new_analyze._selected_images_url)
    '''
    new_analyze.google_search()
    new_summarize = Article("http://www.bbc.com/news/world-asia-43933332")
    new_summarize.download()
    new_summarize.parse()
    new_summarize.text = text
    new_summarize.nlp()
    summarize = new_summarize.summary
    
    print([summarize, new_analyze._selected_images_url])    
    #app.run(debug = True)

    '''